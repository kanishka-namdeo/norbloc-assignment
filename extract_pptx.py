from pptx import Presentation
import sys
import io

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

prs = Presentation(r'd:\test_misc\norbloc-assignment\norbloc_AI_assignment.pptx')

for i, slide in enumerate(prs.slides):
    print(f'\n{"="*80}')
    print(f'SLIDE {i+1}')
    print(f'Layout: {slide.slide_layout.name}')
    print(f'{"="*80}')
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                if para.text.strip():
                    print(f'  {para.text}')
        if shape.has_table:
            print(f'  [TABLE]')
            table = shape.table
            for row_idx, row in enumerate(table.rows):
                row_texts = []
                for cell in row.cells:
                    row_texts.append(cell.text.strip())
                print(f'    Row {row_idx}: {" | ".join(row_texts)}')
        if hasattr(shape, 'image'):
            print(f'  [IMAGE: {shape.shape_id}]')
    print()
