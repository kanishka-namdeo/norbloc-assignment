from pptx import Presentation
import os
import base64

prs = Presentation(r'd:\test_misc\norbloc-assignment\norbloc_AI_assignment.pptx')

img_dir = r'd:\test_misc\norbloc-assignment\pptx_images'
os.makedirs(img_dir, exist_ok=True)

for i, slide in enumerate(prs.slides):
    for j, shape in enumerate(slide.shapes):
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            image = shape.image
            ext = image.ext
            img_data = image.blob
            filename = f'slide{i+1}_img{j}.{ext}'
            filepath = os.path.join(img_dir, filename)
            with open(filepath, 'wb') as f:
                f.write(img_data)
            print(f'Saved: {filepath} ({len(img_data)} bytes)')

# Also extract speaker notes
for i, slide in enumerate(prs.slides):
    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame.text
        if notes.strip():
            print(f'\nSLIDE {i+1} NOTES:')
            print(notes)
